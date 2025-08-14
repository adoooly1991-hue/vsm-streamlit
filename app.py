# app.py (v3)
import streamlit as st
import pandas as pd
import json, io, yaml
from engine import ProcessStep, from_records, compute_lead_time, score_wastes, load_rules, dictify

st.set_page_config(page_title="Automated VSM (v3)", layout="wide")
st.title("Automated Current-State VSM (v3)")
st.caption("Now with material & information flow lanes in PPTX, waste confidence, and PDF export.")

with st.sidebar:
    st.header("1) Load inputs")
    uploaded_csv = st.file_uploader("Upload process CSV", type=["csv"])
    rules_file = st.file_uploader("Custom rules (YAML, optional)", type=["yaml","yml"])
    st.markdown("---")
    st.header("Options")
    available_time_hr = st.number_input("Available time per shift (hours)", value=8.0, min_value=0.1, step=0.5)
    st.markdown("---")
    example_btn = st.button("Use example data", use_container_width=True)

# Example data
if example_btn and not uploaded_csv:
    example_csv = """process_id,process_name,process_type,ct_sec,units_per_period,wip_in_units,defect_pct,rework_pct,downtime_pct,safety_incidents,push_pull,co_freq_per_shift,co_time_min,operators,distance_m,layout_moves,walk_m_per_unit,approval_delays_min,waiting_starved_pct
P1,Receiving,Manual,70,,400,5,13,7,0,Pull,1,10,3,30,1,10,2,5
P2,Inspection,Manual,90,,250,4,6,5,0,Push,2,15,2,10,1,25,0,12
P3,Pick & Pack,Manual,60,,300,6,10,8,1,Pull,3,20,4,120,3,35,10,18
"""
    uploaded_csv = io.BytesIO(example_csv.encode("utf-8"))

# Load CSV and rules
df = None
if uploaded_csv:
    df = pd.read_csv(uploaded_csv)
if rules_file:
    rules = yaml.safe_load(rules_file.read())
else:
    with open("rules.yaml","r") as f:
        rules = yaml.safe_load(f)

if df is not None:
    st.success(f"Loaded {len(df)} process steps")
    st.dataframe(df)

    steps = from_records(df.to_dict(orient="records"))
    st.header("2) Minimal waste questions (only if needed)")
    with st.expander("Open questions"):
        for s in steps:
            with st.container(border=True):
                st.subheader(f"{s.id} – {s.name}")
                c1, c2, c3, c4 = st.columns(4)
                s.answers["d1"] = c1.checkbox("FPY below 97%?", value=(s.defect_pct or 0)>3, key=f"{s.id}_d1")
                s.answers["w1"] = c2.checkbox("Queue time > CT?", value=False, key=f"{s.id}_w1")
                s.answers["op1"] = c3.checkbox("Produce to forecast w/o pull?", value=(s.push_pull=='Push'), key=f"{s.id}_op1")
                s.answers["m1"] = c4.checkbox("Walk > 20 m/unit?", value=(s.walk_m_per_unit or 0)>20, key=f"{s.id}_m1")

    st.header("3) Compute & review")
    available_time_sec = available_time_hr * 3600.0
    result = compute_lead_time(steps, available_time_sec)

    # Waste scoring with confidence
    waste_rows = []
    for s in steps:
        ws = score_wastes(s, rules)
        row = {"process_id": s.id, "process_name": s.name}
        for k,v in ws.scores.items():
            row[f"{k}_score"] = v
            row[f"{k}_confidence"] = ws.confidence.get(k,"Low")
        waste_rows.append(row)
    waste_df = pd.DataFrame(waste_rows)

    c1, c2 = st.columns(2)
    with c1:
        st.metric("Total Lead Time (hh:mm:ss)", pd.to_timedelta(result["lead_time_sec"], unit="s"))
        st.write("CT Bottleneck (sec):", result["ct_bottleneck_sec"])
    with c2:
        st.subheader("Waste scores + confidence")
        st.dataframe(waste_df)

    st.subheader("Per-step timing breakdown (sec)")
    breakdown = pd.DataFrame(result["by_step"]).T.reset_index().rename(columns={"index":"process_id"})
    st.dataframe(breakdown)

    # ---------- PPTX export with lanes ----------
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    def add_process_box(slide, x, y, w, h, title, fields):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(235, 241, 255)
        shape.line.color.rgb = RGBColor(91, 155, 213)
        tf = shape.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True
        for label, val in fields:
            p = tf.add_paragraph()
            p.text = f"{label}: {'' if val is None else val}"
            p.level = 1; p.font.size = Pt(12)

    def draw_arrow(slide, x1, y1, x2, y2, is_pull=False):
        # Simple arrow: line + small triangle
        line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight connector
        line.line.width = Pt(2)
        line.line.color.rgb = RGBColor(0,0,0)
        # Kanban icon for pull
        if is_pull:
            kshape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2-10_000, y2-6_000, 10_000, 6_000)
            kshape.fill.solid(); kshape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            kshape.line.color.rgb = RGBColor(0,0,0)
            kshape.text_frame.text = "KANBAN"

    def build_pptx(steps):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title.text_frame.text = "Current-State VSM"
        title.text_frame.paragraphs[0].font.size = Pt(28)

        # Lanes
        mat_y = Inches(5.2); info_y = Inches(1.2)
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), info_y, Inches(9.2), Inches(0.6)).text_frame.text = "Information Flow"
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), mat_y, Inches(9.2), Inches(0.6)).text_frame.text = "Material Flow"

        # Process boxes horizontally
        start_x = Inches(0.5); gap_x = Inches(3.0); box_w = Inches(2.5); box_h = Inches(2.6)
        box_y = Inches(2.2)
        positions = []
        for i, s in enumerate(steps):
            x = start_x + Inches(i)*gap_x
            positions.append((x, box_y))
            fields = [
                ("Cycle time (sec)", s.ct_sec),
                ("Process type", s.process_type),
                ("Unplanned downtime (%)", s.downtime_pct),
                ("% defects", s.defect_pct),
                ("N. safety issues", s.safety_incidents),
                ("% rework rate", s.rework_pct),
                ("WIP (units)", s.wip_units_in),
                ("Push / Pull", s.push_pull),
                ("Changeover freq/shift", s.co_freq_per_shift),
                ("Changeover time (min)", s.co_time_min),
                ("N. operators", s.operators),
            ]
            add_process_box(slide, x, box_y, box_w, box_h, f"{s.id} – {s.name}", fields)

        # Material arrows between boxes + push/pull icon
        for i in range(len(positions)-1):
            (x1, y1) = positions[i]; (x2, y2) = positions[i+1]
            is_pull = (steps[i].push_pull or "").lower() == "pull"
            draw_arrow(slide, int(x1+box_w), int(y1+box_h/2), int(x2), int(y2+box_h/2), is_pull=is_pull)

        # Information flow (simple: ERP to each step)
        info_src = slide.shapes.add_shape(MSO_SHAPE.FOLDED_CORNER, Inches(0.5), info_y+Inches(0.8), Inches(1.0), Inches(0.8))
        info_src.text_frame.text = "ERP/MRP"
        for (x, y) in positions:
            draw_arrow(slide, int(Inches(1.5)), int(info_y+Inches(1.2)), int(x+box_w/2), int(y), is_pull=False)

        return prs

    st.header("4) Export")
    if st.button("Generate PPTX (lanes + icons)"):
        prs = build_pptx(steps)
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        st.download_button("Download PPTX", data=buf, file_name="vsm_current_state_lanes.pptx")

    # ---------- PDF export replicating data boxes ----------
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import cm
    from reportlab.lib.colors import black

    def build_pdf(steps, result):
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=landscape(A4))
        width, height = landscape(A4)

        # Title
        c.setFont("Helvetica-Bold", 20)
        c.drawString(2*cm, height-1.5*cm, "Current-State VSM (Data Boxes)")

        # Draw process boxes grid (up to 10)
        x0, y0 = 1.5*cm, height-3*cm
        box_w, box_h = 8*cm, 6*cm
        dx = 9.5*cm; dy = 7.5*cm
        for i, s in enumerate(steps[:10]):
            r = i // 3; k = i % 3
            x = x0 + k*dx; y = y0 - r*dy
            c.setLineWidth(1); c.rect(x, y-box_h, box_w, box_h, stroke=1, fill=0)
            c.setFont("Helvetica-Bold", 12)
            c.drawString(x+0.3*cm, y-0.6*cm, f"{s.id} – {s.name}")
            fields = [
                ("Cycle time (sec)", s.ct_sec),
                ("Process type", s.process_type),
                ("Unplanned downtime (%)", s.downtime_pct),
                ("% defects", s.defect_pct),
                ("N. safety issues", s.safety_incidents),
                ("% rework rate", s.rework_pct),
                ("WIP (units)", s.wip_units_in),
                ("Push / Pull", s.push_pull),
                ("Changeover freq/shift", s.co_freq_per_shift),
                ("Changeover time (min)", s.co_time_min),
                ("N. operators", s.operators),
            ]
            c.setFont("Helvetica", 10)
            yy = y-1.2*cm
            for label, val in fields:
                c.drawString(x+0.5*cm, yy, f"{label}: {'' if val is None else val}")
                yy -= 0.45*cm

        # KPI summary
        c.setFont("Helvetica-Bold", 12)
        c.drawString(2*cm, 1.8*cm, f"Total Lead Time: {pd.to_timedelta(result['lead_time_sec'], unit='s')}")
        c.drawString(2*cm, 1.2*cm, f"CT Bottleneck (sec): {result['ct_bottleneck_sec']}")

        c.showPage(); c.save(); buf.seek(0)
        return buf

    if st.button("Generate PDF (template-style)"):
        pdfbuf = build_pdf(steps, result)
        st.download_button("Download PDF", data=pdfbuf, file_name="vsm_current_state.pdf", mime="application/pdf")

    # JSON results
    export_json = {
        "generated_at": pd.Timestamp.utcnow().isoformat(),
        "lead_time_sec": result["lead_time_sec"],
        "ct_bottleneck_sec": result["ct_bottleneck_sec"],
        "by_step": result["by_step"],
        "wastes": waste_df.to_dict(orient="records"),
        "steps": dictify(steps)
    }
    st.download_button("Download results (JSON)", data=json.dumps(export_json, indent=2), file_name="vsm_results.json")

else:
    st.info("Upload a CSV or click 'Use example data' to get started.")
